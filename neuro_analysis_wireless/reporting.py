from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import io



def insert_plot_ppt(pptfile, left, top, width, insertinlastslide=False):
    '''
    input measurements should be in Inches(x)
    '''
    prs = Presentation(pptfile)
    #prs = Presentation()
    #onlytitle_slide_layout = prs.slide_layouts[5]
    blank_slide_layout = prs.slide_layouts[6]
    if insertinlastslide:
        #slide = prs.slides.get(len(prs.slides))
        slide=prs.slides[len(prs.slides)-1]
    else: #INSERT NEW SLIDE WITH MEASUREMEMTS
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        slide = prs.slides.add_slide(blank_slide_layout)
    image_stream = io.BytesIO()
    plt.savefig(image_stream)
    pic = slide.shapes.add_picture(image_stream, left, top, width=width )
    prs.save(pptfile)



def insert_txt_ppt(pptfile, left, top, width, height, stringx, fontsize=28):
    '''
    inputs sould be in Inches(x)
    '''
    prs = Presentation(pptfile)
    slide=prs.slides[len(prs.slides)-1] #get last slide
    #txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox = slide.shapes.add_textbox(left,top,width,height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = stringx
    p.font.size = Pt(fontsize)
    prs.save(pptfile)
    


def insert_shap_summary_ppt(pptfile,shap_values, X_test, column_names, labelsdict, title):

    shap.summary_plot(shap_values, X_test, feature_names=column_names, 
                      class_names=list(labelsdict.keys()), class_inds='original', show=False)
    plt.savefig('tempfig.png')
    plt.clf()

    prs = Presentation(pptfile)
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(blank_slide_layout)

    left=Inches(5)
    top=Inches(1.5)
    height=Inches(7.5)
    _ = slide.shapes.add_picture('tempfig.png', left=left, top=top, height=height )
    prs.save(pptfile)
    insert_txt_ppt(pptfile, Inches(5), 0, width=Inches(10), height=Inches(1), stringx=title,fontsize=20)



def insert_shap_images_by_category_ppt(pptfile, shap_values, X_test, labelsdict, column_names, slidetitle,
                                       imagesperslide=4, slidewidth=16, slideheight=9):
    
    if type(shap_values) is not list: #for binary cases where shap is a single matrix
        shap_values=[shap_values]
        
    counter=0
    for i in range(len(shap_values)):
        #create shap image by category
        shap.summary_plot(shap_values[i], X_test, 
                          title='XXXXXXXX', #list(labelsdict.keys())[0],
                          feature_names=column_names, 
                          class_names=list(labelsdict.keys()),
                          show=False)
        plt.savefig('tempfig.png')
        plt.clf()
        #insert image in ppt in appropriate place, create new slide if needed

        prs = Presentation(pptfile)
        if i%imagesperslide==0: #New slide
            blank_slide_layout = prs.slide_layouts[6]
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            slide = prs.slides.add_slide(blank_slide_layout)
            counter=0
            
            #insert title
            txBox = slide.shapes.add_textbox(left=Inches(1), top=Inches(0.25), width=Inches(10), height=Inches(1))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = slidetitle
            p.font.size = Pt(20)

        else: #use the same slide
            slide=prs.slides[len(prs.slides)-1]
            counter+=1

        left=Inches(counter*(slidewidth/imagesperslide))
        top=Inches(3)
        width=Inches(slidewidth/imagesperslide)
        toptext=top-Inches(1)
        _ = slide.shapes.add_picture('tempfig.png', left=left, top=top, width=width )
        prs.save(pptfile)
        insert_txt_ppt(pptfile, left+Inches(1), toptext, width, height=Inches(1), stringx=list(labelsdict.keys())[i],fontsize=16)
